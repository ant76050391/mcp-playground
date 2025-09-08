"""Office documents extractor (Word, PowerPoint, Excel)."""

import logging
from pathlib import Path
from typing import Optional, Dict, Any, List
import io

from ..core.base import BaseExtractor
from ..core.models import ExtractedDocument, FileType, ChunkInfo

logger = logging.getLogger(__name__)

class OfficeExtractor(BaseExtractor):
    """Extract text from Microsoft Office documents."""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
        self.version = "1.0.0"
        
        # Import libraries and track availability
        self.docx_available = False
        self.pptx_available = False
        self.xlsx_available = False
        
        try:
            from docx import Document
            self.Document = Document
            self.docx_available = True
            logger.debug("python-docx available")
        except ImportError:
            logger.warning("python-docx not available")
        
        try:
            from pptx import Presentation
            self.Presentation = Presentation
            self.pptx_available = True
            logger.debug("python-pptx available")
        except ImportError:
            logger.warning("python-pptx not available")
        
        try:
            from openpyxl import load_workbook
            self.load_workbook = load_workbook
            self.xlsx_available = True
            logger.debug("openpyxl available")
        except ImportError:
            logger.warning("openpyxl not available")
    
    def can_extract(self, file_path: Path) -> bool:
        """Check if this extractor can handle the file."""
        if not file_path.exists() or not file_path.is_file():
            return False
        
        suffix = file_path.suffix.lower()
        
        if suffix in ['.docx', '.doc']:
            return self.docx_available
        elif suffix in ['.pptx', '.ppt']:
            return self.pptx_available
        elif suffix in ['.xlsx', '.xls']:
            return self.xlsx_available
        
        return False
    
    def extract(self, file_path: Path) -> ExtractedDocument:
        """Extract text from Office document."""
        if not self.can_extract(file_path):
            return self._create_failed_result(
                file_path, 
                "Office extractor not available for this file type"
            )
        
        suffix = file_path.suffix.lower()
        
        try:
            if suffix in ['.docx', '.doc']:
                return self._extract_word(file_path)
            elif suffix in ['.pptx', '.ppt']:
                return self._extract_powerpoint(file_path)
            elif suffix in ['.xlsx', '.xls']:
                return self._extract_excel(file_path)
            else:
                return self._create_failed_result(file_path, f"Unsupported file type: {suffix}")
        
        except Exception as e:
            logger.error(f"Failed to extract from {file_path}: {e}")
            return self._create_failed_result(file_path, str(e))
    
    def _extract_word(self, file_path: Path) -> ExtractedDocument:
        """Extract text from Word document."""
        content, extraction_time = self._time_extraction(
            self._extract_word_content, file_path
        )
        
        metadata = self._create_metadata(
            file_path, 
            FileType.DOCX, 
            extraction_time,
            encoding="UTF-8"
        )
        
        # Add Word-specific metadata
        if content:
            metadata.word_count = len(content.split())
            metadata.char_count = len(content)
        
        return ExtractedDocument(
            file_path=str(file_path),
            content=content,
            metadata=metadata,
            success=True
        )
    
    def _extract_word_content(self, file_path: Path) -> str:
        """Extract content from Word document."""
        doc = self.Document(str(file_path))
        content_parts = []
        
        # Extract main document text
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                content_parts.append(text)
        
        # Extract table content
        for table in doc.tables:
            table_content = self._extract_table_content(table)
            if table_content:
                content_parts.append(f"\n{table_content}\n")
        
        return "\n\n".join(content_parts)
    
    def _extract_powerpoint(self, file_path: Path) -> ExtractedDocument:
        """Extract text from PowerPoint presentation."""
        content, extraction_time = self._time_extraction(
            self._extract_pptx_content, file_path
        )
        
        metadata = self._create_metadata(
            file_path, 
            FileType.PPTX, 
            extraction_time,
            encoding="UTF-8"
        )
        
        if content:
            metadata.word_count = len(content.split())
            metadata.char_count = len(content)
            
            # Count slides
            slide_count = content.count("# Slide ")
            if slide_count > 0:
                metadata.page_count = slide_count
        
        return ExtractedDocument(
            file_path=str(file_path),
            content=content,
            metadata=metadata,
            success=True
        )
    
    def _extract_pptx_content(self, file_path: Path) -> str:
        """Extract content from PowerPoint presentation."""
        prs = self.Presentation(str(file_path))
        content_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = [f"# Slide {i}"]
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                
                # Extract table content if present
                if shape.has_table:
                    table_content = self._extract_pptx_table_content(shape.table)
                    if table_content:
                        slide_content.append(table_content)
            
            if len(slide_content) > 1:  # More than just the slide header
                content_parts.append("\n".join(slide_content))
        
        return "\n\n".join(content_parts)
    
    def _extract_excel(self, file_path: Path) -> ExtractedDocument:
        """Extract text from Excel spreadsheet."""
        content, extraction_time = self._time_extraction(
            self._extract_xlsx_content, file_path
        )
        
        metadata = self._create_metadata(
            file_path, 
            FileType.XLSX, 
            extraction_time,
            encoding="UTF-8"
        )
        
        if content:
            metadata.word_count = len(content.split())
            metadata.char_count = len(content)
        
        return ExtractedDocument(
            file_path=str(file_path),
            content=content,
            metadata=metadata,
            success=True
        )
    
    def _extract_xlsx_content(self, file_path: Path) -> str:
        """Extract content from Excel spreadsheet."""
        workbook = self.load_workbook(str(file_path), read_only=True, data_only=True)
        content_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_content = [f"# Sheet: {sheet_name}"]
            
            # Find the actual data range (avoid empty cells)
            rows_with_data = []
            for row in sheet.iter_rows():
                row_data = []
                for cell in row:
                    if cell.value is not None:
                        row_data.append(str(cell.value).strip())
                    else:
                        row_data.append("")
                
                # Only include rows that have some data
                if any(cell_data for cell_data in row_data):
                    rows_with_data.append(" | ".join(row_data))
            
            if rows_with_data:
                sheet_content.extend(rows_with_data)
                content_parts.append("\n".join(sheet_content))
        
        workbook.close()
        return "\n\n".join(content_parts)
    
    def _extract_table_content(self, table) -> str:
        """Extract content from Word table."""
        table_content = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            
            if any(cell_data for cell_data in row_data):
                table_content.append(" | ".join(row_data))
        
        return "\n".join(table_content)
    
    def _extract_pptx_table_content(self, table) -> str:
        """Extract content from PowerPoint table."""
        table_content = []
        
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            
            if any(cell_data for cell_data in row_data):
                table_content.append(" | ".join(row_data))
        
        return "\n".join(table_content)